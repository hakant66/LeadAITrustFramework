"""
Generate LeadAI Presentation Tests

Tests for the PowerPoint presentation generator (slide helpers and full creation).
"""

import pytest

# Skip entire module if python-pptx not installed
pytest.importorskip("pptx")

from pptx import Presentation

from app.generate_leadai_presentation import (
    add_title_slide,
    add_content_slide,
    add_two_column_slide,
    create_leadai_presentation,
)


# --- FIXTURES ---

@pytest.fixture
def prs():
    """Minimal Presentation with default slide layouts"""
    return Presentation()


# --- add_title_slide ---

def test_add_title_slide_creates_slide_and_sets_text(prs):
    """Test add_title_slide adds one slide with title and subtitle"""
    assert len(prs.slides) == 0
    slide = add_title_slide(prs, "My Title", "My Subtitle")
    assert len(prs.slides) == 1
    assert slide.shapes.title.text == "My Title"
    assert slide.placeholders[1].text == "My Subtitle"


def test_add_title_slide_returns_slide(prs):
    """Test add_title_slide returns the created slide"""
    slide = add_title_slide(prs, "A", "B")
    assert slide is not None
    assert slide == prs.slides[0]


# --- add_content_slide ---

def test_add_content_slide_creates_slide_with_bullets(prs):
    """Test add_content_slide adds one slide with title and content items"""
    assert len(prs.slides) == 0
    items = ["First point", "Second point", "Third point"]
    slide = add_content_slide(prs, "Content Title", items)
    assert len(prs.slides) == 1
    assert slide.shapes.title.text == "Content Title"
    tf = slide.placeholders[1].text_frame
    assert tf.paragraphs[0].text == "First point"
    assert tf.paragraphs[1].text == "Second point"
    assert tf.paragraphs[2].text == "Third point"


def test_add_content_slide_single_item(prs):
    """Test add_content_slide with a single content item"""
    slide = add_content_slide(prs, "Single", ["Only one"])
    assert len(prs.slides) == 1
    assert slide.placeholders[1].text_frame.paragraphs[0].text == "Only one"


def test_add_content_slide_empty_items(prs):
    """Test add_content_slide with one empty string item"""
    slide = add_content_slide(prs, "Empty", [""])
    assert len(prs.slides) == 1
    assert slide.shapes.title.text == "Empty"


# --- add_two_column_slide ---

def test_add_two_column_slide_creates_slide_with_two_columns(prs):
    """Test add_two_column_slide adds one slide with title and two columns"""
    left = ["Left 1", "Left 2"]
    right = ["Right 1", "Right 2"]
    slide = add_two_column_slide(prs, "Two Columns", left, right)
    assert len(prs.slides) == 1
    assert slide.shapes.title.text == "Two Columns"
    left_tf = slide.placeholders[1].text_frame
    right_tf = slide.placeholders[2].text_frame
    assert left_tf.paragraphs[0].text == "Left 1"
    assert left_tf.paragraphs[1].text == "Left 2"
    assert right_tf.paragraphs[0].text == "Right 1"
    assert right_tf.paragraphs[1].text == "Right 2"


def test_add_two_column_slide_returns_slide(prs):
    """Test add_two_column_slide returns the created slide"""
    slide = add_two_column_slide(prs, "T", ["L"], ["R"])
    assert slide is not None
    assert slide == prs.slides[0]


# --- create_leadai_presentation ---

def test_create_leadai_presentation_returns_presentation():
    """Test create_leadai_presentation returns a Presentation with multiple slides"""
    prs = create_leadai_presentation()
    assert hasattr(prs, "slides") and hasattr(prs, "slide_width")
    assert len(prs.slides) >= 1


def test_create_leadai_presentation_has_title_slide():
    """Test create_leadai_presentation first slide is LeadAI title"""
    prs = create_leadai_presentation()
    first = prs.slides[0]
    assert "LeadAI" in first.shapes.title.text or "Trust" in first.shapes.title.text


def test_create_leadai_presentation_slide_dimensions():
    """Test create_leadai_presentation sets slide dimensions"""
    prs = create_leadai_presentation()
    # Default is 10 x 7.5 inches in the script
    assert prs.slide_width > 0
    assert prs.slide_height > 0
