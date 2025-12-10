# build_leadai_with_uploaded_logo.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from datetime import date

# Brand tokens (sampled from LeadAI brochure)
LEADAI_YELLOW = RGBColor(0xFF, 0xB4, 0x00)   # #FFB400
LEADAI_DARK_GREY = RGBColor(0x54, 0x58, 0x5A) # #54585A
ACCENT_DARK = RGBColor(0x2D, 0x3E, 0x40)      # #2D3E40

prs = Presentation()
LOGO_FILENAME = "leadai_logo.png"  # put your uploaded logo here

def add_text_logo(slide,left=Inches(0.2),top=Inches(0.15),width=Inches(2.8),height=Inches(0.6),size=28):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    r1 = p.add_run(); r1.text = "Lead"; r1.font.size = Pt(size); r1.font.name="Calibri"; r1.font.color.rgb = LEADAI_DARK_GREY
    r2 = p.add_run(); r2.text = "AI";   r2.font.size = Pt(size); r2.font.name="Calibri"; r2.font.bold=True; r2.font.color.rgb = LEADAI_YELLOW

def add_image_logo(slide,left,top,width=None,height=None):
    if os.path.exists(LOGO_FILENAME):
        try:
            if width is not None:
                slide.shapes.add_picture(LOGO_FILENAME,left,top,width=width)
            elif height is not None:
                slide.shapes.add_picture(LOGO_FILENAME,left,top,height=height)
            else:
                slide.shapes.add_picture(LOGO_FILENAME,left,top)
            return True
        except Exception as e:
            print("Warning: couldn't add image logo:", e)
            return False
    return False

def add_diagonal_accent(slide):
    left = prs.slide_width - Inches(1.6)
    top = -Inches(0.4)
    width = Inches(2.5)
    height = Inches(1.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LEADAI_YELLOW
    shape.line.fill.background()
    shape.rotation = -20

def add_bullet_slide(title, bullets, accent=True):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_DARK

    logo_added = False
    if os.path.exists(LOGO_FILENAME):
        logo_added = add_image_logo(slide, left=Inches(0.2), top=Inches(0.12), width=Inches(1.8))
    if not logo_added:
        add_text_logo(slide,left=Inches(0.2),top=Inches(0.12),width=Inches(2.0),height=Inches(0.4),size=18)

    if accent:
        add_diagonal_accent(slide)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    first = True
    for line in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.font.color.rgb = ACCENT_DARK
    return slide

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "AI Governance & EU AI Act Readiness"
subtitle = slide.placeholders[1]
subtitle.text = f"Leadership Briefing – {date.today().strftime('%d %B %Y')}"
subtitle.text_frame.paragraphs[0].font.size = Pt(16)
subtitle.text_frame.paragraphs[0].font.color.rgb = LEADAI_DARK_GREY

if os.path.exists(LOGO_FILENAME):
    try:
        add_image_logo(slide,left=Inches(0.3),top=Inches(0.3),width=prs.slide_width - Inches(1.0))
    except Exception:
        add_text_logo(slide,left=Inches(0.3),top=Inches(0.3),width=Inches(9.0),height=Inches(1.0),size=36)
else:
    add_text_logo(slide,left=Inches(0.3),top=Inches(0.3),width=Inches(9.0),height=Inches(1.0),size=36)
add_diagonal_accent(slide)

# (Add all content slides — same as previous versions)
add_bullet_slide("Why AI Governance Matters Now", [
    "AI is now embedded in core business processes and customer journeys.",
    "Regulators expect systematic governance, not ad-hoc control.",
    "EU AI Act introduces strict obligations for high-risk and general-purpose AI.",
    "ISO/IEC 42001 provides a structured AI Management System (AIMS).",
    "Proactive governance reduces risk while enabling safe, scalable innovation."
])
# ... (repeat other content slides exactly as earlier)
add_bullet_slide("Key Findings – Current State", [
    "AI governance foundations exist, but enterprise AI ownership and accountability remain informal.",
    "AI-specific risk assessments are performed inconsistently across teams.",
    "Data governance is solid for structured data, but fairness and bias controls are not yet standardised.",
    "Technical capabilities are competent, yet monitoring, drift detection, and robustness testing vary.",
    "Third-party AI components lack structured due diligence and compliance documentation."
])
add_bullet_slide("Proposed Roadmap (Aligned to Reality)", [
    "0–3 months: Establish AI governance council; define enterprise AI policy; complete AI system inventory and EU AI Act classification.",
    "3–6 months: Implement foundational AIMS processes (risk, lifecycle, data, documentation); publish standards and templates.",
    "6–12 months: Deploy monitoring, drift detection, bias testing, incident reporting, and human-oversight controls for high-risk systems.",
    "12+ months: Scale AIMS maturity across teams; integrate governance workflows with MLOps; prepare for internal/external audits.",
    "Ongoing: Review metrics, incidents, regulatory updates; refine controls as AI adoption expands."
])

# Contact slide
contact_slide = prs.slides.add_slide(prs.slide_layouts[1])
contact_slide.shapes.title.text = "Contact"
if os.path.exists(LOGO_FILENAME):
    add_image_logo(contact_slide,left=Inches(0.25),top=Inches(0.2),width=Inches(3.0))
else:
    add_text_logo(contact_slide,left=Inches(0.25),top=Inches(0.2),width=Inches(3.0),height=Inches(0.6),size=20)
body = contact_slide.shapes.placeholders[1].text_frame
body.clear()
p = body.paragraphs[0]
p.text = "theleadai.co.uk\n+44 782 499 8608\n+44 739 356 5252\ncan@theleadai.co.uk"
p.font.size = Pt(16)
p.font.color.rgb = LEADAI_DARK_GREY

# Final "Thank you" slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_diagonal_accent(slide)
if os.path.exists(LOGO_FILENAME):
    add_image_logo(slide,left=Inches(0.2),top=Inches(0.12),width=Inches(1.8))
else:
    add_text_logo(slide,left=Inches(0.2),top=Inches(0.12),width=Inches(2.0),height=Inches(0.4),size=18)

tx = slide.shapes.add_textbox(Inches(1),Inches(2.2),prs.slide_width - Inches(2),Inches(1.5))
tf = tx.text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
r = p.add_run()
r.text = "Thank you"
r.font.size = Pt(44); r.font.bold = True; r.font.name = "Calibri"; r.font.color.rgb = ACCENT_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
r2 = p2.add_run()
r2.text = "Questions? contact@theleadai.co.uk"
r2.font.size = Pt(18); r2.font.name = "Calibri"; r2.font.color.rgb = LEADAI_DARK_GREY

output_name = "AI_Governance_Leadership_Deck_LeadAI_Branded_Thanks_withLogo.pptx"
prs.save(output_name)
print(f"Saved: {output_name}")
