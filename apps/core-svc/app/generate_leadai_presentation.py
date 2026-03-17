#!/usr/bin/env python3
"""
Generate LeadAI Process Presentation as PowerPoint

Creates a PowerPoint presentation (.pptx) describing the LeadAI Trust Framework
end-to-end governance process.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is not installed.")
    print("Install it with: pip install python-pptx")
    exit(1)

from datetime import datetime


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide


def add_content_slide(prs, title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = content_shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
    
    return slide


def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content layout
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Left column
    left_shape = slide.placeholders[1]
    tf_left = left_shape.text_frame
    tf_left.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = item
        p.level = 0
    
    # Right column
    right_shape = slide.placeholders[2]
    tf_right = right_shape.text_frame
    tf_right.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = item
        p.level = 0
    
    return slide


def create_leadai_presentation():
    """Create the LeadAI process presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "LeadAI Trust Framework",
        "Automated AI Governance & Compliance Platform\nFrom Regulations to Trust Scores"
    )
    
    # Slide 2: Executive Summary
    add_two_column_slide(
        prs,
        "Executive Summary",
        [
            "The Challenge",
            "• Multiple Regulations",
            "• Complex Requirements",
            "• Evidence Management",
            "• Continuous Monitoring"
        ],
        [
            "The Solution",
            "• Automated Governance",
            "• Evidence-Centric",
            "• Continuous Scoring",
            "• Intelligent Reporting"
        ]
    )
    
    # Slide 3: Process Overview
    add_content_slide(
        prs,
        "Process Overview - The Big Picture",
        [
            "1. Regulatory Input Funnel → EU AI Act, ISO 42001, NIST AI RMF",
            "2. Policy & Guardrail Engine → Policies, Guardrails, KPIs, Controls",
            "3. Evidence Collection → Manual uploads + Automatic integrations (Jira)",
            "4. Evidence Vault → Centralized storage with integrity verification",
            "5. Control & Audit Process → Human review + AI assistance",
            "6. Artifact & Scoring Engine → Trust scores (TOL-0 to TOL-3)",
            "7. Monitoring & Reporting → Dashboards + Executive reports"
        ]
    )
    
    # Slide 4: Phase 1 - Regulatory Input Funnel
    add_content_slide(
        prs,
        "Phase 1: Regulatory Input Funnel",
        [
            "EU AI Act:",
            "  • High-Risk AI Systems (Annex III)",
            "  • Transparency Requirements (Article 50)",
            "  • Data Governance (Article 10)",
            "",
            "ISO/IEC 42001:",
            "  • Organizational Context (Clause 4)",
            "  • Planning & Operation (Clauses 6-8)",
            "  • Performance Evaluation (Clause 9)",
            "",
            "NIST AI RMF:",
            "  • Govern, Map, Measure, Manage"
        ]
    )
    
    # Slide 5: Phase 2 - Policy & Guardrail Engine
    add_content_slide(
        prs,
        "Phase 2: Policy & Guardrail Engine",
        [
            "Policy Generation:",
            "  • Automatic mapping: Requirements → Policies",
            "  • Policy templates and customization",
            "",
            "Guardrail Rules:",
            "  • YAML-driven rule definitions",
            "  • Hard gates (non-negotiable)",
            "  • Soft gates (warnings)",
            "",
            "KPI/Metric Definition:",
            "  • KPI registry with normalization",
            "  • Control mapping to pillars",
            "  • Evidence requirements"
        ]
    )
    
    # Slide 6: Phase 3 - Evidence Collection
    add_two_column_slide(
        prs,
        "Phase 3: Evidence Collection Layer",
        [
            "Manual Collection:",
            "• Direct file uploads",
            "• Metadata entry",
            "• Control linking",
            "• Status tracking"
        ],
        [
            "Automatic Integration:",
            "• Jira sync (requirements, risks, tests)",
            "• REST API endpoints",
            "• Scheduled syncs",
            "• Webhook support"
        ]
    )
    
    # Slide 7: Phase 4 - Evidence Vault
    add_content_slide(
        prs,
        "Phase 4: Evidence Vault",
        [
            "Storage:",
            "  • Centralized repository (S3-compatible)",
            "  • Version control",
            "  • Retention policies",
            "",
            "Integrity Verification:",
            "  • SHA-256 hashing",
            "  • Hash verification",
            "  • Digital signatures",
            "",
            "Metadata Management:",
            "  • Rich metadata (project, control, status)",
            "  • Search & discovery",
            "  • Relationship linking"
        ]
    )
    
    # Slide 8: Phase 5 - Control & Audit
    add_two_column_slide(
        prs,
        "Phase 5: Control & Audit Process",
        [
            "AI-Assisted Review:",
            "• Automated analysis",
            "• Compliance checking",
            "• Anomaly detection",
            "• Recommendations"
        ],
        [
            "Human Decision:",
            "• Evidence verification",
            "• Challenge process",
            "• Approval workflow",
            "• Rejection handling"
        ]
    )
    
    # Slide 9: Phase 6 - Scoring
    add_content_slide(
        prs,
        "Phase 6: Artifact & Scoring Engine",
        [
            "Trust Scoring:",
            "  • Trust Axes: Safety, Compliance, Provenance",
            "  • Trust Verdict: TOL-0 (Critical) to TOL-3 (Compliant)",
            "  • Provenance Levels: P0 (Fail) to P3 (Excellent)",
            "",
            "Pillar Scoring:",
            "  • Governance (GOV)",
            "  • Trust, Control & Transparency (TCT)",
            "  • Data Governance (DG)",
            "  • Risk Management (RM)",
            "  • Provenance (PROV)",
            "",
            "Score Calculation:",
            "  KPI → Control → Pillar → Overall Score"
        ]
    )
    
    # Slide 10: Phase 7 - Monitoring & Reporting
    add_two_column_slide(
        prs,
        "Phase 7: Monitoring & Reporting",
        [
            "Real-Time Dashboard:",
            "• Project-level view",
            "• Portfolio-level view",
            "• Trust scorecards",
            "• KPI trends",
            "• Alert summary"
        ],
        [
            "Executive Reporting:",
            "• Template-based reports",
            "• AI-generated summaries",
            "• Scheduled delivery",
            "• Multiple formats (PDF, Excel, PPT)"
        ]
    )
    
    # Slide 11: Recursive Process
    add_content_slide(
        prs,
        "Recursive Process - Continuous Improvement",
        [
            "The process is recursive and continuous:",
            "",
            "1. Regulatory Updates → New/updated requirements",
            "2. Policy Refresh → Update policies and guardrails",
            "3. Evidence Re-collection → Gather new evidence",
            "4. Re-audit & Re-scoring → Verify and recalculate",
            "5. Monitoring & Alerts → Track changes and anomalies",
            "6. Reporting & Feedback → Generate insights",
            "",
            "Frequency:",
            "  • Daily: Evidence collection, score updates",
            "  • Weekly: Evidence audit, compliance checks",
            "  • Monthly: Policy review, KPI assessment",
            "  • Quarterly: Regulatory updates",
            "  • Annually: Comprehensive audit"
        ]
    )
    
    # Slide 12: Key Benefits
    add_content_slide(
        prs,
        "Key Benefits & Value Proposition",
        [
            "For Compliance Teams:",
            "  ✅ Automated mapping: Regulations → Policies → KPIs",
            "  ✅ Evidence management: Centralized, verified",
            "  ✅ Audit trail: Complete traceability",
            "",
            "For Risk Management:",
            "  ✅ Risk visibility: Comprehensive risk register",
            "  ✅ Risk scoring: Quantified risk levels",
            "  ✅ Trend analysis: Risk evolution tracking",
            "",
            "For Executives:",
            "  ✅ Executive dashboards: High-level visibility",
            "  ✅ Automated reporting: AI-powered insights",
            "  ✅ Trust scores: Simple, actionable metrics"
        ]
    )
    
    # Slide 13: Conclusion
    add_content_slide(
        prs,
        "Conclusion",
        [
            "LeadAI Trust Framework Delivers:",
            "",
            "✅ Automated Governance: From regulations to policies to KPIs",
            "✅ Evidence Management: Centralized, verified, traceable",
            "✅ Continuous Scoring: Real-time trust and compliance scores",
            "✅ Intelligent Reporting: AI-powered executive insights",
            "✅ Recursive Process: Continuous improvement and monitoring",
            "",
            "Key Differentiators:",
            "  • Regulatory-Agnostic: Supports multiple frameworks",
            "  • Evidence-Centric: Built around evidence verification",
            "  • Human + AI: Combines AI assistance with human judgment",
            "  • Recursive: Self-improving, continuous process"
        ]
    )
    
    return prs


if __name__ == '__main__':
    print("Generating LeadAI Process Presentation...")
    prs = create_leadai_presentation()
    output_path = 'docs/LeadAI_Process_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
